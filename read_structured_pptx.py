import json
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER


def extract_text_frame(shape):
    data = []

    for paragraph in shape.text_frame.paragraphs:
        text = paragraph.text.strip()
        if not text:
            continue

        data.append({
            "text": text,
            "level": paragraph.level
        })

    return data


def identify_shape_type(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "GROUP"

    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type

        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
            return "TITLE_PLACEHOLDER"

        if ph_type == PP_PLACEHOLDER.BODY:
            return "BODY_PLACEHOLDER"

        return f"PLACEHOLDER_{ph_type}"

    if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        return "TEXTBOX"

    if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
        return "AUTO_SHAPE"

    return str(shape.shape_type)


def extract_shape(shape):
    shape_data = {
        "type": identify_shape_type(shape),
        "left": shape.left,
        "top": shape.top,
        "width": shape.width,
        "height": shape.height,
    }

    if shape.has_text_frame:
        shape_data["text_content"] = extract_text_frame(shape)

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        shape_data["children"] = [
            extract_shape(subshape) for subshape in shape.shapes
        ]

    return shape_data


def extract_presentation(path):
    prs = Presentation(path)

    presentation_data = {
        "slides": []
    }

    for slide_index, slide in enumerate(prs.slides):
        slide_data = {
            "slide_index": slide_index,
            "layout_name": slide.slide_layout.name,
            "shapes": []
        }

        # Extract title separately if exists
        if slide.shapes.title:
            slide_data["title"] = slide.shapes.title.text.strip()

        for shape in slide.shapes:
            slide_data["shapes"].append(extract_shape(shape))

        presentation_data["slides"].append(slide_data)

    return presentation_data


def format_shape_reading_order(shape_data, element_index, slide_num, indent=0):
    """Format a single shape's reading order information as text lines.

    Reading order for python-pptx: This is the raw iteration order of shapes
    on each slide, which corresponds to the XML element order in the slide's
    shape tree (spTree). This is the drawing order (back-to-front z-order),
    which may differ from visual reading order. Group shapes are recursively
    expanded with their children listed in their internal order.
    """
    prefix = "  " * indent
    lines = []
    shape_type = shape_data.get("type", "UNKNOWN")
    left = shape_data.get("left")
    top = shape_data.get("top")
    width = shape_data.get("width")
    height = shape_data.get("height")

    # Convert EMUs to approximate cm for readability (1 cm = 360000 EMU)
    pos_str = ""
    if left is not None and top is not None:
        pos_str = f"pos=({left}, {top}) size=({width}x{height}) EMU"

    lines.append(f"{prefix}[Slide {slide_num}] Element #{element_index} | "
                 f"Type: {shape_type} | {pos_str}")

    # Print text content if available
    text_content = shape_data.get("text_content", [])
    for tc in text_content:
        text = tc.get("text", "")
        level = tc.get("level", 0)
        lines.append(f"{prefix}  Text (level {level}): {text}")

    # Recurse into group children
    children = shape_data.get("children", [])
    for child_idx, child in enumerate(children):
        lines.append(f"{prefix}  --- Group child #{child_idx} ---")
        lines.extend(format_shape_reading_order(child, child_idx, slide_num, indent + 2))

    return lines


if __name__ == "__main__":
    file_path = "test-with-images.pptx"  # change path if needed

    structured_data = extract_presentation(file_path)

    with open("pptx-structure/test_with_images_structure.json", "w", encoding="utf-8") as f:
        json.dump(structured_data, f, ensure_ascii=False, indent=2)

    print("Extraction complete. JSON saved as presentation_index.json")