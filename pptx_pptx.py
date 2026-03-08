import os
import sys
import glob as globmod
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def print_shape_info(shape, slide_num, elem_index, indent=0):
    """Format shape info lines for reading order output.

    Reading order for python-pptx (basic): Same as read_structured_pptx.py — iterates
    shapes in XML spTree order (z-order). This lightweight version prints
    shape name, type, and basic position without deep text extraction.
    """
    prefix = "  " * indent
    lines = []

    shape_type = str(shape.shape_type).replace("MSO_SHAPE_TYPE.", "")
    name = shape.name or "(unnamed)"
    has_text = shape.has_text_frame
    text_preview = ""
    if has_text:
        text_preview = shape.text_frame.text.strip()[:120]

    lines.append(f"{prefix}[Slide {slide_num}] Element #{elem_index} | "
                 f"Name: {name} | Type: {shape_type} | "
                 f"pos=({shape.left}, {shape.top}) size=({shape.width}x{shape.height}) EMU")

    if text_preview:
        lines.append(f"{prefix}  Text: {text_preview}")

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_idx, child in enumerate(shape.shapes):
            lines.append(f"{prefix}  --- Group child #{child_idx} ---")
            lines.extend(print_shape_info(child, slide_num, child_idx, indent + 2))

    return lines


def extract_reading_order(pptx_path, output_dir="reading_order_outputs"):
    """Extract and save reading order using basic python-pptx iteration.

    Reading order: Shapes are iterated in their XML shape-tree order (spTree),
    which is the drawing z-order (back-to-front). This is python-pptx's native
    iteration order and reflects how PowerPoint stores shapes internally.
    """
    os.makedirs(output_dir, exist_ok=True)

    prs = Presentation(pptx_path)
    pptx_name = os.path.splitext(os.path.basename(pptx_path))[0]
    output_filename = f"pptx__{pptx_name}.txt"
    output_path = os.path.join(output_dir, output_filename)

    lines = [
        f"Reading Order Report — python-pptx",
        f"Source file: {pptx_path}",
        f"Library: python-pptx (lightweight iteration)",
        f"Order semantics: Raw shape-tree iteration order (XML spTree z-order).",
        "=" * 80,
    ]

    for slide_index, slide in enumerate(prs.slides):
        slide_num = slide_index + 1
        layout_name = slide.slide_layout.name
        title_text = slide.shapes.title.text.strip() if slide.shapes.title else ""

        lines.append("")
        lines.append(f"--- Slide {slide_num} (layout: {layout_name}) ---")
        if title_text:
            lines.append(f"    Title: {title_text}")
        lines.append("")

        for elem_idx, shape in enumerate(slide.shapes):
            lines.extend(print_shape_info(shape, slide_num, elem_idx))
            lines.append("")

    text = "\n".join(lines)
    with open(output_path, "w", encoding="utf-8") as f:
        f.write(text)

    print(f"[python-pptx basic] Reading order saved to: {output_path}")
    return output_path

def classify_text(shape):
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                font_size = run.font.size.pt
                if font_size > 35:
                    return "TITLE"
    return "TEXT"


if __name__ == "__main__":
    if "--reading" in sys.argv:
        pptx_files = globmod.glob("*.pptx")
        if not pptx_files:
            print("No PPTX files found in current directory.")
        for pf in pptx_files:
            extract_reading_order(pf)
    else:
        prs = Presentation("test-with-groups.pptx")
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                print(shape.name)
                text += shape.name + ": "
                if shape.has_text_frame:
                    label = classify_text(shape)
                    print(f"[{label}] {shape.text_frame.text}")
                    text += f"[{label}] {shape.text_frame.text}\n"
                elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    text += "\n"
                    for shape in shape.shapes:
                        print(shape.name)
                        text += "   " + shape.name + ": "
                        if shape.has_text_frame:
                            label = classify_text(shape)
                            print(f"[{label}] {shape.text_frame.text}")
                            text += f"[{label}] {shape.text_frame.text}\n"
                else:
                    text += "\n"

        with open("pptx-parser/test-with-groups.txt", "w", encoding="utf-8") as f:
            f.write(text)